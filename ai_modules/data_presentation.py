import collections 
import collections.abc
from pptx import Presentation
import PIL
import os
import mysql.connector
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
import numpy as np


import warnings
warnings.filterwarnings("ignore")
import pandas as pd
import os

import sys
from os.path import exists
import math
import random
from dotenv import load_dotenv
from pathlib import Path

dotenv_path = Path(os.path.abspath(__file__) + '../.env')
load_dotenv()

def typerin(v):
    s = str(v)
    regex = re.compile(r'(?P<list>\[[^]]+\])|(?P<float>\d*\.\d+)|(?P<int>\d+)|(?P<string>[a-zA-Z]+)')
    return  regex.search(s).lastgroup

def execute(connection, query, cursor):
    if (connection['type'] == 'csv' or connection['type'] == 'json') and df is not None:
      import pandasql as ps
      import json
      cursor = ps.sqldf(query, locals()).to_json(orient="records")
      cursor = json.loads(cursor)
      fks = None
    else:
      cursor.execute(query)
    if connection['type'] == 'mysql':
      cursor = cursor.fetchall()
    res = []
    for row in cursor:
      for col in row.keys():
        if not row[col] == None and (not (type(row[col]) == int or type(row[col]) == float)):
            row[col] = str(row[col])
      res.append(row)
    return res

def drawChart(data, shapes):

    stat = data.describe()
    total_length = len(data.index)
    relevant = []
    unrelevant = []
    import datetime

    countryInfo = pd.read_json(
    'https://raw.githubusercontent.com/lukes/ISO-3166-Countries-with-Regional-Codes/master/all/all.json'
    )

    countryNames = list(countryInfo['name'])
    countryA2 = list(countryInfo['alpha-2'])
    countryA3 = list(countryInfo['alpha-3'])

    mapColumn = None
    dateColumn = None
    for col in data:
        uniques = pd.unique(data[col])
        if data[col][0] in countryNames or data[col][0] in countryA2 or data[col][0] in countryA3:
            mapColumn = col
        if len(uniques) > 0.70 * total_length:
            relevant.append(col)
        if len(list(filter(lambda x: str(x) not in str([None, 'nan', 'NaN', 0]), uniques))) == 0:
            unrelevant.append(col)
        if isinstance(data[col][0] , datetime.date) or isinstance(data[col][0] , datetime.datetime):
            dateColumn = col

    baseW = Cm(20)
    baseH = Cm(12)

    restat = {
    "std": None
    }
    try:
        restat = stat.idxmax(axis=1)
    except:
        restat['std'] = None

    import numbers
    import altair as alt
    alt.renderers.enable('altair_saver', fmts=['vega-lite', 'png'])
    firstRow = data.iloc[0]
    from pptx.chart.data import CategoryChartData, ChartData
    from pptx.enum.chart import XL_CHART_TYPE


    if mapColumn is not None and total_length > 0:
        from vega_datasets import data as vegaDT
        
        relevantId = 0
        if restat['std'] is not None:
            value = restat['std']
        else:
            while relevantId < len(relevant) - 1 and mapColumn == relevant[relevantId]:
                relevantId = relevantId + 1
            if mapColumn == relevant[relevantId]:
                relevantId = 0
                keys = firstRow.keys()
                while relevantId < len(keys) - 1 and mapColumn == keys[relevantId]:
                    relevantId = relevantId + 1
                value = keys[relevantId]
            else:
                value = relevant[relevantId]

        value = value.replace('(','').replace(')','')
        geoData = data
        geoData.columns = geoData.columns.str.replace(r"[()]", "")
        geoData['CODE'] = 0
        for index, row in geoData.iterrows():
            x = row[mapColumn]
            #pycountry.countries
            country = countryInfo[(countryInfo['name'] == x) | (countryInfo['alpha-2'] == x) | (countryInfo['alpha-3'] == x)]
            if len(country['country-code'].values) > 0:
                geoData.at[index, 'CODE'] = country['alpha-3'].values[0]
            geoData.at[index, 'id'] = int(index)

        df = geoData[['id', 'CODE', value]]

        import geopandas
        from matplotlib import pyplot as plt

        world = geopandas.read_file(geopandas.datasets.get_path('naturalearth_lowres'))

        world.columns=['pop_est', 'continent', 'name', 'CODE', 'gdp_md_est', 'geometry']

        world['value'] = 0
        for index, row in world.iterrows():
            rr = df[df['CODE'] == row['CODE']]

            if rr is not None and len(rr[value].values) > 0:
                world.at[index, 'value'] = np.sum(rr[value].values)

        world = world[world['CODE'] != 'ATA']

        fig = plt.figure()
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis('off')
        world.plot(ax = ax, column='value',
            figsize=(250, 200),cmap='terrain')
        
        
        parcell = ''.join(random.choice('0123456789ABCDEF') for i in range(16))
        name = os.path.dirname(__file__) + "\\..\\temp\\" + parcell + '.png'
        plt.savefig(name, transparent=True, dpi=1200)
        subtitle = shapes[0]
        subtitle.text = "Let's take a look at our world"
        chart = shapes.add_picture(name, 0, 0, Cm(23.57), Cm(20.37))
        os.remove(name)
        return [chart,Cm(23.57), Cm(20.37)]
    elif dateColumn is not None:
        if restat['std'] is not None:
            chart_data = ChartData()
            chart_data.categories = data[dateColumn]
            chart_data.add_series(restat['std'].capitalize().replace('_',' '), list(data[restat['std']]))

            chart = shapes.add_chart(
                XL_CHART_TYPE.AREA,  0, 0, baseW, baseH, chart_data
            ).chart
            chart.series[0].smooth = True
            return [chart,baseW, baseH]
        elif len(relevant) > 0 and isinstance(firstRow[relevant[0]], numbers.Number):
            chart_data = ChartData()
            chart_data.categories = data[dateColumn]
            chart_data.add_series(relevant[0].capitalize().replace('_',' '), list(data[relevant[0]]))

            chart = shapes.add_chart(
                XL_CHART_TYPE.AREA,  0, 0, baseW, baseH, chart_data
            ).chart
            chart.series[0].smooth = True
            return [chart,baseW, baseH]
        elif len(relevant) > 0:
            chart_data = CategoryChartData()
            chart_data.categories = data[dateColumn]
            chart_data.add_series(relevant[0].capitalize().replace('_',' '), list(data[relevant[0]]))
            return [shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, baseW, baseH, chart_data
            ).chart,baseW, baseH]
    else:
        if restat['std'] is not None:
            relevantId = 0
            while relevantId < len(relevant) - 1 and restat['std'] == relevant[relevantId]:
                relevantId = relevantId + 1
            if (len(relevant) > 0 and isinstance(firstRow[relevant[relevantId]], numbers.Number)):
                chart_data = CategoryChartData()
                chart_data.categories = data[restat['std']]
                chart_data.add_series(relevant[relevantId].capitalize().replace('_',' '), list(data[relevant[relevantId]]))
                chart =  shapes.add_chart(
                    XL_CHART_TYPE.AREA, 0, 0, baseW, baseH, chart_data
                ).chart
                chart.series[0].smooth = True
                return [chart,baseW, baseH]
            # quantitative + quantitative = line
            elif len(relevant) > 0:
                from pptx.chart.data import CategoryChartData
                from pptx.enum.chart import XL_CHART_TYPE
                chart_data = CategoryChartData()
                chart_data.categories = data[relevant[relevantId]]
                chart_data.add_series(restat['std'].capitalize().replace('_',' '), list(data[restat['std']]))
                if random.random() > 0.5:
                    return [shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, baseW, baseH, chart_data
                    ).chart,baseW, baseH]
                else:
                    return [shapes.add_chart(
                        XL_CHART_TYPE.DOUGHNUT, 0, 0, baseW, baseH, chart_data
                    ).chart,baseW, baseH]
            # text + quantitative = bar or pie
        elif len(relevant) > 0:
            if (len(relevant) > 1 and typerin(firstRow[relevant[1]]) in ['float','int']):
                chart_data = ChartData()
                chart_data.categories = data[relevant[0]]
                chart_data.add_series(relevant[1].capitalize().replace('_',' '), list(data[relevant[1]]))
                return [shapes.add_chart(
                    XL_CHART_TYPE.AREA, 0, 0, baseW, baseH, chart_data
                ).chart,baseW, baseH]
            elif len(relevant) > 1  and typerin(firstRow[relevant[0]]) in ['float','int']:
                
                chart_data = CategoryChartData()
                chart_data.categories = data[relevant[1]]
                chart_data.add_series(relevant[0].capitalize().replace('_',' '), list(data[relevant[0]]))
                if random.random() > 0.5:
                    return [shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, baseW, baseH, chart_data
                    ).chart,baseW, baseH]
                else:
                    return [shapes.add_chart(
                        XL_CHART_TYPE.DOUGHNUT, 0, 0, baseW, baseH, chart_data
                    ).chart,baseW, baseH]

prs=Presentation('template.pptx')

db = sys.argv[1]
if (db == None):
  raise Exception('No db')

import json

queries = sys.argv[2]
if (queries == None):
  raise Exception('No query')
else:
    queries = json.loads(queries)

titles = sys.argv[3]
if (titles == None):
  raise Exception('No query')
else:
    titles = json.loads(titles)
    if len(titles) != len(queries):
        raise Exception('Misalignment')

slides = {
    "header": 0,
    "subheader": 1,
    "text": 2,
    "tablex": 3,
    "chart": 4,
    "blank": 5,
    "singleton": 6,
    "table": 7
}


mydb = mysql.connector.connect(
  host=os.getenv("CODEX_DB_HOST"),
  user=os.getenv("CODEX_DB_USER"),
  password=os.getenv("CODEX_DB_PASS"),
  database=os.getenv("CODEX_DB_NAME")
)

mycursor = mydb.cursor(dictionary=True)


import pymssql

mycursor.execute("SELECT * FROM `db` WHERE id = %s", (int(db),))

connection = mycursor.fetchone()


if connection['type'] == 'mssql':
  conn = pymssql.connect(connection['server'], connection['username'], connection['password'], connection['name'])
  execution = conn.cursor(as_dict=True)
elif connection['type'] == 'mysql':
  conn = mysql.connector.connect(
    host=connection['server'],
    user=connection['username'],
    password=connection['password'],
    database=connection['name']
  )
  execution = conn.cursor(dictionary=True)


slide = prs.slides.add_slide(prs.slide_layouts[slides["header"]])


texts = '\n'.join(titles)
import openai
openai.api_key = os.getenv("GPT_CODEX_API_KEY")

gpt_prompt = '''Rewrite each prompt as a capturing introduction title.
Reply using this format, without repeating the original input: <index>|<prompt>
''' + texts

response = openai.ChatCompletion.create(
model="gpt-3.5-turbo",
messages=[{
    "role": "user", "content": gpt_prompt
}],
temperature=0.3,
max_tokens=150,
top_p=1.0,
frequency_penalty=0.0,
presence_penalty=0.0,
stop=["#", ";"]
)
response = response['choices'][0]['message']['content']



parts = response.split('\n')
parts = list(filter(lambda x: x.strip() != '', parts))

import re
from datetime import datetime

gpt_prompt = '''Create an introduction title given these prompts, without repeating the original input
''' + texts

response = openai.ChatCompletion.create(
model="gpt-3.5-turbo",
messages=[{
    "role": "user", "content": gpt_prompt
}],
temperature=0.2,
max_tokens=150,
top_p=1.0,
frequency_penalty=0.0,
presence_penalty=0.0,
stop=["#", ";"]
)
response = response['choices'][0]['message']['content']

titleText = re.sub(r'Introduction Title:', '', response.replace('"','')).strip()

title=slide.shapes.title # assigning a title
title.text=titleText #"Let's explore data" # title
slide.shapes[1].text = datetime.today().strftime('%Y-%m-%d')



idPrompt = 0
for query in queries:

    copiedHeader = prs.slides.add_slide(prs.slide_layouts[slides["subheader"]])

    subText = re.split('\?|!|\.', parts[idPrompt].split('|')[1])
    subText = list(filter(lambda x: x.strip() != '', subText))
    copiedHeader.shapes.title.text = subText[0] #prompt["prompt"].capitalize().replace('_',' ')
    if len(subText) > 0:
        copiedHeader.shapes[1].text = '. '.join(subText[1:]).strip() #prompt["prompt"].capitalize().replace('_',' ')
    idPrompt = idPrompt + 1
    df = pd.DataFrame(execute({'type':connection['type']}, query, execution))
    if df.shape[0] > 1:
        copiedTable = prs.slides.add_slide(prs.slide_layouts[slides["table"]])
        subtitle = copiedTable.shapes[0]

        shape = copiedTable.shapes.add_table(min(df.shape[0] + 1, 11), df.shape[1], 0, 0, Cm(23), Cm(10.5))
        shape.left = int((prs.slide_width / 2) - (shape.width / 2))
        shape.top = int((prs.slide_height / 2) - (shape.height / 2) - Cm(1))
        table = shape.table

        idxC = 0
        for col in df:
            cell = table.cell(0, idxC)
            cell.text = str(col).capitalize().replace('_',' ')
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.font.name  = 'HP Simplified Hans'
            p.font.bold = True
            p.alignmnet = PP_ALIGN.CENTER
            idxC = idxC + 1
        for index, row in df.iterrows():
            if index > 9:
                subtitle.text = 'But many more are present'
                break
            idxC = 0
            for col in row:
                cell = table.cell(index + 1, idxC)
                p = cell.text_frame.paragraphs[0]
                p.font.name  = 'HP Simplified Hans'
                cell.text = str(col).capitalize().replace('_',' ')
                idxC = idxC + 1

        copiedImg = prs.slides.add_slide(prs.slide_layouts[slides["table"]])
        # copiedImg.shapes[3].text = subText[0]
        [chart, width, height] = drawChart(df, copiedImg.shapes)
        try:
            chart.font.name = 'HP Simplified Hans'
            chart.value_axis.tick_labels.font.name = 'HP Simplified Hans'
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.tick_labels.font.size = Pt(7)
            chart.category_axis.tick_labels.font.name = 'HP Simplified Hans'
            chart.category_axis.has_major_gridlines = False
            chart.category_axis.tick_labels.font.size = Pt(7)
            chart.has_legend = False 
            chart.series[0].smooth = True
        except:
            chart = chart

        try:
            chart.data_labels.show_category_name = True
            chart.legends.font.name = 'HP Simplified Hans'
        except:
            chart = chart
        copiedImg.shapes[1].left = int((prs.slide_width / 2) - (width / 2))
        copiedImg.shapes[1].top = int((prs.slide_height / 2) - (height / 2))
    else:
        copiedTable = prs.slides.add_slide(prs.slide_layouts[slides["singleton"]])
        import numbers
        if  isinstance(df.iloc[0,0], numbers.Number):
            copiedTable.shapes[0].text = "Just" if df.iloc[0,0] < 3 else "About"
        else:
            copiedTable.shapes[0].text = "It's"
        copiedTable.shapes[1].text = str(df.iloc[0,0])
        if len(df.iloc[0]) > 1:            
            copiedTable.shapes[2].text = 'With ' + ', '.join(str(x) for x in dict(df.iloc[0,1:]).values())


ending = prs.slides.add_slide(prs.slide_layouts[slides["header"]])
ending.shapes.title.text = "And... that's it for the moment"
ending.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
ending.shapes[1].text = "Autogenerated by Queric exploiting GPT"
xml_slides = prs.slides._sldIdLst  
slides = list(xml_slides)
for idx in range(0, 7):
    xml_slides.remove(slides[idx]) 

name = ''.join(random.choice('0123456789ABCDEF') for i in range(16))
prs.save(os.path.dirname(__file__) + "\\..\\temp\\" + name + ".pptx") # saving file

print(name + ".pptx")
